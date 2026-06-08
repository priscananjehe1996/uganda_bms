import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import json
import shutil
from pathlib import Path

# Setup directories
SOURCE_DIR = Path(r"D:\OneDrive\Bridge stuff")
DEST_DIR = Path(r"D:\OneDrive\Bridge stuff\uganda_bms\public\data\extracted_photos")
OUTPUT_JSON = Path(r"D:\OneDrive\Bridge stuff\uganda_bms\public\data\extracted_metadata.json")

# Ensure clean destination directory
if DEST_DIR.exists():
    shutil.rmtree(DEST_DIR)
os.makedirs(DEST_DIR, exist_ok=True)

# List of excluded dirs
EXCLUDED_DIRS = ['node_modules', '.git', 'uganda_bms', 'uganda_bms_backend', 'uganda_bms_deploy', 'traffic-spatial-worktree', '.tmp.driveupload']

extracted_metadata = []

def extract_from_pdf(filepath):
    text_snippet = ""
    try:
        doc = fitz.open(filepath)
        for i, page in enumerate(doc):
            text_snippet += page.get_text() + " "
            if len(text_snippet) > 2000:
                break
            
            # Optionally extract images from PDF (first 5 images to save time)
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list[:5]):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                photo_name = f"{filepath.stem}_page{i}_img{img_index}.{image_ext}"
                photo_path = DEST_DIR / photo_name
                with open(photo_path, "wb") as f:
                    f.write(image_bytes)
                
                extracted_metadata.append({
                    "parent_file": filepath.name,
                    "type": "PHOTO",
                    "filename": photo_name,
                    "filepath": str(photo_path)
                })
        
        doc.close()
    except Exception as e:
        print(f"Failed to read PDF {filepath.name}: {e}")
    return text_snippet[:1000]

def extract_from_docx(filepath):
    text_snippet = ""
    try:
        doc = docx.Document(filepath)
        text_snippet = " ".join([p.text for p in doc.paragraphs])[:1000]
        
        # Extract images from docx
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                photo_name = f"{filepath.stem}_{os.path.basename(rel.target_ref)}"
                photo_path = DEST_DIR / photo_name
                with open(photo_path, "wb") as f:
                    f.write(rel.target_part.blob)
                    
                extracted_metadata.append({
                    "parent_file": filepath.name,
                    "type": "PHOTO",
                    "filename": photo_name,
                    "filepath": str(photo_path)
                })
                
    except Exception as e:
        print(f"Failed to read DOCX {filepath.name}: {e}")
    return text_snippet

def extract_from_pptx(filepath):
    text_snippet = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_snippet += shape.text + " "
                if shape.shape_type == 13: # Picture
                    image = shape.image
                    photo_name = f"{filepath.stem}_{image.filename}"
                    photo_path = DEST_DIR / photo_name
                    with open(photo_path, "wb") as f:
                        f.write(image.blob)
                        
                    extracted_metadata.append({
                        "parent_file": filepath.name,
                        "type": "PHOTO",
                        "filename": photo_name,
                        "filepath": str(photo_path)
                    })
    except Exception as e:
        print(f"Failed to read PPTX {filepath.name}: {e}")
    return text_snippet[:1000]

def process_directory(directory):
    for root, dirs, files in os.walk(directory):
        dirs[:] = [d for d in dirs if d not in EXCLUDED_DIRS]
        for file in files:
            # Skip temp files
            if file.startswith('~'):
                continue
                
            filepath = Path(root) / file
            size_mb = os.path.getsize(filepath) / (1024 * 1024)
            
            # Skip massive files to save time
            if size_mb > 100:
                continue
            
            ext = filepath.suffix.lower()
            snippet = ""
            file_type = "UNKNOWN"
            
            if ext == ".pdf":
                file_type = "PDF"
                snippet = extract_from_pdf(filepath)
            elif ext == ".docx":
                file_type = "DOCX"
                snippet = extract_from_docx(filepath)
            elif ext == ".pptx":
                file_type = "PPTX"
                snippet = extract_from_pptx(filepath)
            elif ext in [".txt", ".csv"]:
                file_type = ext.upper().replace('.', '')
                try:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        snippet = f.read(1000)
                except:
                    pass
            elif ext in [".png", ".jpg", ".jpeg"]:
                file_type = "PHOTO"
                
            if file_type != "UNKNOWN":
                extracted_metadata.append({
                    "filename": file,
                    "filepath": str(filepath),
                    "type": file_type,
                    "snippet": snippet.strip().replace('\n', ' ') if snippet else "",
                    "size_mb": round(size_mb, 2)
                })

process_directory(SOURCE_DIR)

with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
    json.dump(extracted_metadata, f, indent=2)

print(f"Extraction complete! Found {len(extracted_metadata)} items. Wrote to {OUTPUT_JSON}")
